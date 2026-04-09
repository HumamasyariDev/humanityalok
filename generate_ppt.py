"""
SmartPark PPT Presentation Generator
Generates a professional 10-slide presentation about the SmartPark project.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── Color palette ──
BLUE_DARK   = RGBColor(0x1E, 0x3A, 0x5F)
BLUE_MED    = RGBColor(0x2B, 0x6C, 0xB3)
BLUE_LIGHT  = RGBColor(0x3B, 0x82, 0xF6)
BLUE_BG     = RGBColor(0xDB, 0xEA, 0xFE)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x00, 0x00, 0x00)
GRAY_DARK   = RGBColor(0x33, 0x33, 0x33)
GRAY_MED    = RGBColor(0x66, 0x66, 0x66)
GRAY_LIGHT  = RGBColor(0xF1, 0xF5, 0xF9)
GRAY_BG     = RGBColor(0xF8, 0xFA, 0xFC)
GREEN       = RGBColor(0x10, 0xB9, 0x81)
ORANGE      = RGBColor(0xF5, 0x9E, 0x0B)
RED         = RGBColor(0xEF, 0x44, 0x44)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SW = prs.slide_width
SH = prs.slide_height


# ── Helpers ──

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_gradient_rect(slide, left, top, width, height, color1, color2):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    fill = shape.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = color2
    fill.gradient_stops[1].position = 1.0
    return shape

def add_rect(slide, left, top, w, h, fill_color=None, border_color=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, left, top, w, h)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def set_text(shape, text, size=14, color=GRAY_DARK, bold=False, align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf

def add_text_box(slide, left, top, w, h, text, size=14, color=GRAY_DARK, bold=False, align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def add_multiline(slide, left, top, w, h, lines, size=13, color=GRAY_DARK, line_spacing=1.3, font_name="Segoe UI", bullet=False):
    """lines = list of (text, bold, color_override)"""
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, bld, clr = item, False, color
        else:
            txt = item[0]
            bld = item[1] if len(item) > 1 else False
            clr = item[2] if len(item) > 2 else color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(size)
        p.font.color.rgb = clr
        p.font.bold = bld
        p.font.name = font_name
        p.space_after = Pt(4)
        if bullet and not bld:
            p.level = 0
    return txBox

def add_card(slide, left, top, w, h, title, body_lines, accent_color=BLUE_MED, title_size=14, body_size=11):
    card = add_rect(slide, left, top, w, h, fill_color=WHITE, border_color=RGBColor(0xE2,0xE8,0xF0))
    # Accent top bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent_color
    bar.line.fill.background()
    # Title
    add_text_box(slide, left + Inches(0.2), top + Inches(0.12), w - Inches(0.4), Inches(0.35),
                 title, size=title_size, color=BLUE_DARK, bold=True)
    # Body
    add_multiline(slide, left + Inches(0.2), top + Inches(0.5), w - Inches(0.4), h - Inches(0.6),
                  body_lines, size=body_size, color=GRAY_MED)

# ── Flowchart helpers ──

def fc_oval(slide, cx, cy, w, h, text, fill=WHITE, stroke=BLUE_DARK):
    left = cx - w // 2
    top = cy - h // 2
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = stroke
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.color.rgb = BLUE_DARK
    p.font.bold = True
    p.font.name = "Segoe UI"
    shape.text_frame.paragraphs[0].space_before = Pt(0)
    shape.text_frame.paragraphs[0].space_after = Pt(0)
    return shape

def fc_box(slide, cx, cy, w, h, text, fill=WHITE, stroke=BLUE_MED, font_size=8):
    left = cx - w // 2
    top = cy - h // 2
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = stroke
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = GRAY_DARK
    p.font.name = "Segoe UI"
    return shape

def fc_diamond(slide, cx, cy, w, h, text, fill=WHITE, stroke=ORANGE):
    left = cx - w // 2
    top = cy - h // 2
    shape = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = stroke
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(7)
    p.font.color.rgb = GRAY_DARK
    p.font.name = "Segoe UI"
    return shape

def fc_arrow(slide, x1, y1, x2, y2, color=BLUE_MED):
    """Simple line arrow from (x1,y1) to (x2,y2)"""
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # 1 = straight
    connector.line.color.rgb = color
    connector.line.width = Pt(1.2)
    # Add arrowhead via oxml
    spPr = connector._element.find(qn('a:ln'))
    if spPr is None:
        cxnSp = connector._element
        spPr_el = cxnSp.find(qn('p:spPr'))
        if spPr_el is None:
            spPr_el = cxnSp.find(qn('p:cxnSp'))
        ln = spPr_el.find(qn('a:ln')) if spPr_el is not None else None
    else:
        ln = spPr
    # Fallback: access via _element tree
    cxnSp = connector._element
    spPr_parent = cxnSp.find(qn('p:spPr'))
    if spPr_parent is not None:
        ln_el = spPr_parent.find(qn('a:ln'))
        if ln_el is not None:
            tail = ln_el.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
            ln_el.append(tail)
    return connector

def fc_arrow_down(slide, cx, y_from, y_to, color=BLUE_MED):
    return fc_arrow(slide, cx, y_from, cx, y_to, color)


# ══════════════════════════════════════════════════
# SLIDE 1: COVER
# ══════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_gradient_rect(sl, 0, 0, SW, SH, BLUE_DARK, RGBColor(0x0F, 0x25, 0x40))

# Decorative shapes
deco1 = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-1.5), Inches(5), Inches(5))
deco1.fill.solid()
deco1.fill.fore_color.rgb = RGBColor(0x2B, 0x4F, 0x7F)
deco1.line.fill.background()
deco1.rotation = 15.0

deco2 = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(4.5), Inches(4), Inches(4))
deco2.fill.solid()
deco2.fill.fore_color.rgb = RGBColor(0x1A, 0x30, 0x50)
deco2.line.fill.background()

# Title cluster
add_text_box(sl, Inches(1.5), Inches(1.5), Inches(10), Inches(0.6),
             "SMARTPARK", size=18, color=BLUE_LIGHT, bold=True)
add_text_box(sl, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
             "Sistem Manajemen Parkir\nTerintegrasi AI", size=40, color=WHITE, bold=True)
# Divider line
line = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(3.4), Inches(2.5), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = BLUE_LIGHT
line.line.fill.background()

add_text_box(sl, Inches(1.5), Inches(3.7), Inches(8), Inches(0.8),
             "Aplikasi berbasis web untuk pengelolaan operasional parkir dengan\nfitur AI Chatbot berbasis NVIDIA API", 
             size=16, color=RGBColor(0xA0, 0xC4, 0xF0))

add_text_box(sl, Inches(1.5), Inches(5.5), Inches(8), Inches(0.7),
             "Laravel 12  |  React 19  |  Vite  |  Tailwind CSS  |  NVIDIA AI API",
             size=14, color=RGBColor(0x80, 0xA0, 0xC0))

add_text_box(sl, Inches(1.5), Inches(6.2), Inches(8), Inches(0.5),
             "HumamasyariDev  -  2026",
             size=13, color=RGBColor(0x70, 0x90, 0xB0))


# ══════════════════════════════════════════════════
# SLIDE 2: DESKRIPSI PROGRAM
# ══════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, GRAY_BG)
# Top bar
add_rect(sl, 0, 0, SW, Inches(0.9), fill_color=BLUE_DARK)
add_text_box(sl, Inches(0.8), Inches(0.15), Inches(10), Inches(0.6),
             "01  |  Deskripsi Program", size=24, color=WHITE, bold=True)

# Main content box
add_rect(sl, Inches(0.8), Inches(1.3), Inches(11.7), Inches(2.5), fill_color=WHITE, border_color=RGBColor(0xE2,0xE8,0xF0))
add_multiline(sl, Inches(1.1), Inches(1.5), Inches(11.2), Inches(2.2), [
    ("SmartPark adalah sistem informasi manajemen parkir berbasis web yang dikembangkan", False, GRAY_DARK),
    ("menggunakan arsitektur pemisahan frontend dan backend (decoupled architecture).", False, GRAY_DARK),
    ("", False),
    ("Backend dibangun dengan Laravel 12 + Sanctum Authentication, sedangkan frontend", False, GRAY_DARK),
    ("menggunakan React 19 + Vite + Tailwind CSS. Terintegrasi dengan AI Chatbot berbasis", False, GRAY_DARK),
    ("NVIDIA API (LLaMA 3.1 & Phi-3.5 Vision) untuk cek biaya & status parkir secara real-time.", False, GRAY_DARK),
], size=14, color=GRAY_DARK)

# Feature cards
feats = [
    ("Transaksi Parkir", "Pencatatan masuk/keluar\nkendaraan, perhitungan\ntarif otomatis per jam", BLUE_MED),
    ("Barcode System", "Kode parkir PKR-XXXXXX\ndengan barcode, scan\nuntuk proses keluar", GREEN),
    ("AI Chatbot", "Cek biaya & status via\nteks, kamera, atau\nupload gambar struk", ORANGE),
    ("Multi-Role Access", "3 role: Admin (full),\nPetugas (operasional),\nOwner (monitoring)", RGBColor(0x8B, 0x5C, 0xF6)),
    ("Dashboard & Rekap", "Statistik real-time,\ngrafik 7 hari, rekap\npendapatan, export CSV", RED),
]
card_w = Inches(2.15)
card_h = Inches(1.8)
start_x = Inches(0.8)
gap = Inches(0.2)
for i, (title, body, accent) in enumerate(feats):
    x = start_x + i * (card_w + gap)
    y = Inches(4.2)
    add_card(sl, x, y, card_w, card_h, title, [body], accent_color=accent, body_size=11)


# ══════════════════════════════════════════════════
# SLIDE 3: ARSITEKTUR SISTEM
# ══════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, GRAY_BG)
add_rect(sl, 0, 0, SW, Inches(0.9), fill_color=BLUE_DARK)
add_text_box(sl, Inches(0.8), Inches(0.15), Inches(10), Inches(0.6),
             "02  |  Arsitektur Sistem", size=24, color=WHITE, bold=True)

# Frontend box
add_rect(sl, Inches(0.6), Inches(1.3), Inches(5.5), Inches(5.5), fill_color=RGBColor(0xEF,0xF6,0xFF), border_color=BLUE_LIGHT)
add_text_box(sl, Inches(0.8), Inches(1.4), Inches(5), Inches(0.4),
             "FRONTEND  (React 19 + Vite + Tailwind CSS v4)", size=13, color=BLUE_MED, bold=True)

fe_items = [
    "Pages (11 halaman):",
    "  Login, Dashboard, Users, TarifParkir,",
    "  AreaParkir, Kendaraan, TransaksiMasuk,",
    "  TransaksiKeluar, TransaksiList, Rekap, LogAktivitas",
    "",
    "Components (3):",
    "  Layout.jsx  -  Sidebar + Navbar wrapper",
    "  ChatbotWidget.jsx  -  AI Chatbot (kamera, upload, barcode)",
    "  ConfirmDialog.jsx  -  Modal konfirmasi",
    "",
    "State Management:",
    "  AuthContext.jsx  -  Token, user, role checking",
    "  Axios interceptor  -  Auto Bearer token + 401 redirect",
]
add_multiline(sl, Inches(0.9), Inches(1.85), Inches(5.0), Inches(4.5), fe_items, size=11, color=GRAY_DARK)

# Backend box
add_rect(sl, Inches(6.5), Inches(1.3), Inches(6.2), Inches(5.5), fill_color=RGBColor(0xF0,0xFD,0xF4), border_color=GREEN)
add_text_box(sl, Inches(6.7), Inches(1.4), Inches(5.5), Inches(0.4),
             "BACKEND  (Laravel 12 + Sanctum + MySQL)", size=13, color=GREEN, bold=True)

be_items = [
    "Controllers (8):",
    "  AuthController  -  Login, Logout, Me",
    "  UserController  -  CRUD pengguna",
    "  TarifParkirController  -  CRUD tarif",
    "  AreaParkirController  -  CRUD area parkir",
    "  KendaraanController  -  CRUD kendaraan",
    "  TransaksiController  -  Masuk, Keluar, Struk, Barcode, Rekap, Dashboard",
    "  LogAktivitasController  -  Riwayat aktivitas",
    "  ChatbotController  -  AI Chat + Scan Image (NVIDIA API)",
    "",
    "Models (6):  User, Transaksi, Kendaraan, TarifParkir, AreaParkir, LogAktivitas",
    "Middleware:  RoleMiddleware  -  Otorisasi role-based (admin, petugas, owner)",
    "Database:  MySQL (parkir_ukk) - 6 tabel + personal_access_tokens",
]
add_multiline(sl, Inches(6.8), Inches(1.85), Inches(5.7), Inches(4.5), be_items, size=11, color=GRAY_DARK)


# ══════════════════════════════════════════════════
# SLIDE 4: STRUKTUR DATABASE
# ══════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, GRAY_BG)
add_rect(sl, 0, 0, SW, Inches(0.9), fill_color=BLUE_DARK)
add_text_box(sl, Inches(0.8), Inches(0.15), Inches(10), Inches(0.6),
             "03  |  Struktur Basis Data", size=24, color=WHITE, bold=True)

# Table cards - 6 tables in 2 rows of 3
tables = [
    ("tb_user", "id_user  (PK)\nnama_lengkap  varchar(50)\nusername  varchar(50) UNIQUE\npassword  varchar(100)\nrole  enum(admin, petugas, owner)\nstatus_aktif  tinyint"),
    ("tb_kendaraan", "id_kendaraan  (PK)\nplat_nomor  varchar(15) UNIQUE\njenis_kendaraan  varchar(20)\nwarna  varchar(20)\npemilik  varchar(100)\nid_user  (FK -> tb_user)"),
    ("tb_tarif", "id_tarif  (PK)\njenis_kendaraan  enum(motor, mobil, lainnya)\ntarif_per_jam  decimal(10,0)"),
    ("tb_area_parkir", "id_area  (PK)\nnama_area  varchar(50)\nkapasitas  int\nterisi  int  (default: 0)"),
    ("tb_transaksi", "id_parkir  (PK)  =  Kode PKR\nid_kendaraan  (FK)\nwaktu_masuk / waktu_keluar  datetime\nid_tarif (FK) / id_area (FK) / id_user (FK)\ndurasi_jam  int / biaya_total  decimal\nstatus  enum(masuk, keluar)"),
    ("tb_log_aktivitas", "id_log  (PK)\nid_user  (FK -> tb_user)\naktivitas  varchar(100)\nwaktu_aktivitas  datetime"),
]

card_w = Inches(3.85)
card_h_row1 = Inches(2.5)
card_h_row2 = Inches(2.5)
for i, (tname, tbody) in enumerate(tables):
    row = i // 3
    col = i % 3
    x = Inches(0.6) + col * (card_w + Inches(0.15))
    y = Inches(1.2) + row * (Inches(2.7))
    ch = card_h_row1 if row == 0 else card_h_row2
    accent = BLUE_MED if i in [0, 4] else GREEN if i in [1, 3] else ORANGE
    add_card(sl, x, y, card_w, ch, tname, [tbody], accent_color=accent, title_size=13, body_size=10)

# Relasi note
add_text_box(sl, Inches(0.6), Inches(6.6), Inches(12), Inches(0.5),
             "Relasi:  tb_user -> tb_transaksi (1:N)  |  tb_kendaraan -> tb_transaksi (1:N)  |  tb_tarif -> tb_transaksi (1:N)  |  tb_area_parkir -> tb_transaksi (1:N)  |  tb_user -> tb_log_aktivitas (1:N)",
             size=11, color=GRAY_MED)


# ══════════════════════════════════════════════════
# SLIDE 5: FITUR PER ROLE
# ══════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, GRAY_BG)
add_rect(sl, 0, 0, SW, Inches(0.9), fill_color=BLUE_DARK)
add_text_box(sl, Inches(0.8), Inches(0.15), Inches(10), Inches(0.6),
             "04  |  Hak Akses & Fitur Per Role", size=24, color=WHITE, bold=True)

# Admin column
col_w = Inches(3.8)
col_h = Inches(5.3)

# Admin
add_rect(sl, Inches(0.6), Inches(1.2), col_w, col_h, fill_color=WHITE, border_color=BLUE_MED)
bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(1.2), col_w, Inches(0.5))
bar.fill.solid()
bar.fill.fore_color.rgb = BLUE_MED
bar.line.fill.background()
add_text_box(sl, Inches(0.8), Inches(1.22), col_w, Inches(0.45), "ADMIN (Full Access)", size=16, color=WHITE, bold=True)
admin_feats = [
    "Dashboard  -  Ringkasan statistik",
    "Kelola User  -  CRUD admin/petugas/owner",
    "Tarif Parkir  -  CRUD tarif per jenis",
    "Area Parkir  -  CRUD area + kapasitas",
    "Kelola Kendaraan  -  CRUD data kendaraan",
    "Transaksi Masuk  -  Input kendaraan masuk",
    "Transaksi Keluar  -  Scan barcode + bayar",
    "Daftar Transaksi  -  Riwayat + cetak struk",
    "Rekap Transaksi  -  Grafik + tabel + CSV",
    "Log Aktivitas  -  Riwayat semua aksi",
    "AI Chatbot  -  Cek biaya/status via AI",
]
add_multiline(sl, Inches(0.8), Inches(1.85), Inches(3.4), Inches(4.5), admin_feats, size=11, color=GRAY_DARK)

# Petugas
x2 = Inches(4.65)
add_rect(sl, x2, Inches(1.2), col_w, col_h, fill_color=WHITE, border_color=GREEN)
bar2 = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x2, Inches(1.2), col_w, Inches(0.5))
bar2.fill.solid()
bar2.fill.fore_color.rgb = GREEN
bar2.line.fill.background()
add_text_box(sl, x2 + Inches(0.2), Inches(1.22), col_w, Inches(0.45), "PETUGAS (Operasional)", size=16, color=WHITE, bold=True)
petugas_feats = [
    "Dashboard  -  Ringkasan statistik",
    "Kelola Kendaraan  -  CRUD data kendaraan",
    "Transaksi Masuk  -  Input kendaraan masuk",
    "Transaksi Keluar  -  Scan barcode + bayar",
    "Daftar Transaksi  -  Riwayat + cetak struk",
    "AI Chatbot  -  Cek biaya/status via AI",
]
add_multiline(sl, x2 + Inches(0.2), Inches(1.85), Inches(3.4), Inches(4.5), petugas_feats, size=11, color=GRAY_DARK)

# Owner
x3 = Inches(8.7)
add_rect(sl, x3, Inches(1.2), col_w, col_h, fill_color=WHITE, border_color=ORANGE)
bar3 = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x3, Inches(1.2), col_w, Inches(0.5))
bar3.fill.solid()
bar3.fill.fore_color.rgb = ORANGE
bar3.line.fill.background()
add_text_box(sl, x3 + Inches(0.2), Inches(1.22), col_w, Inches(0.45), "OWNER (Monitoring)", size=16, color=WHITE, bold=True)
owner_feats = [
    "Dashboard  -  Ringkasan statistik",
    "Rekap Transaksi  -  Grafik pendapatan",
    "  per hari, tabel transaksi, ringkasan",
    "  per jenis kendaraan, export CSV",
    "AI Chatbot  -  Cek biaya/status via AI",
]
add_multiline(sl, x3 + Inches(0.2), Inches(1.85), Inches(3.4), Inches(4.5), owner_feats, size=11, color=GRAY_DARK)


# ══════════════════════════════════════════════════
# SLIDE 5: FLOW ALUR KERJA - TRANSAKSI PARKIR (UTAMA)
# ══════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, GRAY_BG)
add_rect(sl, 0, 0, SW, Inches(0.9), fill_color=BLUE_DARK)
add_text_box(sl, Inches(0.8), Inches(0.15), Inches(10), Inches(0.6),
             "05  |  Flow Alur Kerja  -  Transaksi Parkir", size=24, color=WHITE, bold=True)

# Flowchart area
fc_bg = add_rect(sl, Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.9), fill_color=WHITE, border_color=RGBColor(0xE2,0xE8,0xF0))

# === MASUK (Left side) ===
add_text_box(sl, Inches(0.8), Inches(1.15), Inches(5), Inches(0.35),
             "ALUR KENDARAAN MASUK", size=13, color=BLUE_MED, bold=True)

bw = Inches(1.8)
bh = Inches(0.48)
ow = Inches(1.3)
oh = Inches(0.45)
dw = Inches(1.6)
dh = Inches(0.7)

cx_left = Inches(3.2)
y_start = Inches(1.8)
y_gap = Inches(0.75)

# Mulai
fc_oval(sl, cx_left, y_start, ow, oh, "Mulai", fill=BLUE_BG, stroke=BLUE_MED)
y = y_start + Inches(0.22)
fc_arrow_down(sl, cx_left, y + Inches(0.0), y + Inches(0.28))

# Petugas input plat
y2 = y_start + y_gap
fc_box(sl, cx_left, y2, bw, bh, "Petugas Input\nPlat Nomor & Jenis", fill=GRAY_LIGHT)
fc_arrow_down(sl, cx_left, y2 + bh//2, y2 + bh//2 + Inches(0.28))

# Auto-create kendaraan
y3 = y2 + y_gap
fc_box(sl, cx_left, y3, bw, bh, "Sistem Auto-Create\nKendaraan (firstOrCreate)", fill=GRAY_LIGHT)
fc_arrow_down(sl, cx_left, y3 + bh//2, y3 + bh//2 + Inches(0.28))

# Cek area
y4 = y3 + y_gap
fc_diamond(sl, cx_left, y4, dw, dh, "Area Parkir\nPenuh?", fill=RGBColor(0xFF,0xF7,0xED), stroke=ORANGE)

# Ya -> Tolak
fc_box(sl, cx_left + Inches(2.2), y4, bw*3//4, bh*3//4, "Tolak!\nArea Penuh", fill=RGBColor(0xFE,0xE2,0xE2), stroke=RED, font_size=7)
fc_arrow(sl, cx_left + dw//2, y4, cx_left + Inches(2.2) - bw*3//8, y4, ORANGE)
add_text_box(sl, cx_left + Inches(0.8), y4 - Inches(0.22), Inches(0.4), Inches(0.2), "Ya", size=8, color=RED, bold=True)

# Tidak -> lanjut
fc_arrow_down(sl, cx_left, y4 + dh//2, y4 + dh//2 + Inches(0.28))
add_text_box(sl, cx_left + Inches(0.1), y4 + Inches(0.25), Inches(0.5), Inches(0.2), "Tidak", size=8, color=GREEN, bold=True)

# Simpan transaksi
y5 = y4 + y_gap + Inches(0.1)
fc_box(sl, cx_left, y5, bw, bh, "Simpan Transaksi\nStatus = 'masuk'", fill=RGBColor(0xEC,0xFD,0xF5), stroke=GREEN)
fc_arrow_down(sl, cx_left, y5 + bh//2, y5 + bh//2 + Inches(0.28))

# Cetak barcode
y6 = y5 + y_gap
fc_box(sl, cx_left, y6, bw, bh, "Cetak Kartu Parkir\nBarcode PKR-XXXXXX", fill=RGBColor(0xEC,0xFD,0xF5), stroke=GREEN)
fc_arrow_down(sl, cx_left, y6 + bh//2, y6 + bh//2 + Inches(0.28))

# Selesai
y7 = y6 + y_gap
fc_oval(sl, cx_left, y7, ow, oh, "Selesai", fill=BLUE_BG, stroke=BLUE_MED)


# === KELUAR (Right side) ===
add_text_box(sl, Inches(7.0), Inches(1.15), Inches(5), Inches(0.35),
             "ALUR KENDARAAN KELUAR", size=13, color=GREEN, bold=True)

cx_right = Inches(9.5)

# Mulai
fc_oval(sl, cx_right, y_start, ow, oh, "Mulai", fill=BLUE_BG, stroke=BLUE_MED)
fc_arrow_down(sl, cx_right, y_start + oh//2, y_start + oh//2 + Inches(0.28))

# Scan barcode
y2r = y_start + y_gap
fc_box(sl, cx_right, y2r, bw, bh, "Petugas Scan Barcode\natau Input Kode PKR", fill=GRAY_LIGHT)
fc_arrow_down(sl, cx_right, y2r + bh//2, y2r + bh//2 + Inches(0.28))

# Cari transaksi
y3r = y2r + y_gap
fc_box(sl, cx_right, y3r, bw, bh, "Sistem Cari Transaksi\n(status = 'masuk')", fill=GRAY_LIGHT)
fc_arrow_down(sl, cx_right, y3r + bh//2, y3r + bh//2 + Inches(0.28))

# Hitung biaya
y4r = y3r + y_gap
fc_box(sl, cx_right, y4r, bw + Inches(0.3), bh + Inches(0.1), "Hitung Biaya\nceil(menit/60) x tarif_per_jam", fill=RGBColor(0xFF,0xF7,0xED), stroke=ORANGE)
fc_arrow_down(sl, cx_right, y4r + (bh + Inches(0.1))//2, y4r + (bh + Inches(0.1))//2 + Inches(0.28))

# Proses pembayaran
y5r = y4r + y_gap + Inches(0.1)
fc_box(sl, cx_right, y5r, bw, bh, "Proses Pembayaran\n+ Update Status 'keluar'", fill=RGBColor(0xEC,0xFD,0xF5), stroke=GREEN)
fc_arrow_down(sl, cx_right, y5r + bh//2, y5r + bh//2 + Inches(0.28))

# Cetak struk
y6r = y5r + y_gap
fc_box(sl, cx_right, y6r, bw, bh, "Cetak Struk Parkir\n(Detail + Barcode)", fill=RGBColor(0xEC,0xFD,0xF5), stroke=GREEN)
fc_arrow_down(sl, cx_right, y6r + bh//2, y6r + bh//2 + Inches(0.28))

# Selesai
y7r = y6r + y_gap
fc_oval(sl, cx_right, y7r, ow, oh, "Selesai", fill=BLUE_BG, stroke=BLUE_MED)

# Divider
div = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.3), Inches(1.3), Pt(2), Inches(5.5))
div.fill.solid()
div.fill.fore_color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
div.line.fill.background()


# ══════════════════════════════════════════════════
# SLIDE 8: FLOW ALUR KERJA PER ROLE
# ══════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, GRAY_BG)
add_rect(sl, 0, 0, SW, Inches(0.9), fill_color=BLUE_DARK)
add_text_box(sl, Inches(0.8), Inches(0.15), Inches(10), Inches(0.6),
             "06  |  Flow Alur Kerja Per Role", size=24, color=WHITE, bold=True)

# 3 columns for 3 roles
col_w = Inches(3.9)

def draw_role_flow(sl, x_start, y_start, role_name, color, steps):
    """Draw a simplified vertical flowchart for a role"""
    cx = x_start + col_w // 2
    # Role title
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x_start, y_start, col_w, Inches(0.4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    add_text_box(sl, x_start + Inches(0.1), y_start + Inches(0.02), col_w - Inches(0.2), Inches(0.35),
                 role_name, size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    bw_role = Inches(3.0)
    bh_role = Inches(0.36)
    oh_role = Inches(0.34)
    y = y_start + Inches(0.65)
    
    for i, (step_type, text) in enumerate(steps):
        if step_type == 'oval':
            fc_oval(sl, cx, y, Inches(1.1), oh_role, text, fill=BLUE_BG, stroke=color)
            if i < len(steps) - 1:
                fc_arrow_down(sl, cx, y + oh_role//2, y + oh_role//2 + Inches(0.17), color)
            y += Inches(0.52)
        elif step_type == 'box':
            fc_box(sl, cx, y, bw_role, bh_role, text, fill=GRAY_LIGHT, stroke=color, font_size=8)
            if i < len(steps) - 1:
                fc_arrow_down(sl, cx, y + bh_role//2, y + bh_role//2 + Inches(0.17), color)
            y += Inches(0.52)
        elif step_type == 'diamond':
            dh_role = Inches(0.52)
            fc_diamond(sl, cx, y, Inches(1.8), dh_role, text, fill=RGBColor(0xFF,0xF7,0xED), stroke=color)
            if i < len(steps) - 1:
                fc_arrow_down(sl, cx, y + dh_role//2, y + dh_role//2 + Inches(0.17), color)
            y += Inches(0.68)

# Admin flow
admin_steps = [
    ('oval', 'Login Admin'),
    ('box', 'Dashboard Admin'),
    ('diamond', 'Pilih Menu'),
    ('box', 'Kelola User/Tarif/Area/Kendaraan'),
    ('box', 'Transaksi Masuk/Keluar'),
    ('box', 'Rekap + Log Aktivitas'),
    ('box', 'AI Chatbot'),
    ('oval', 'Logout'),
]
draw_role_flow(sl, Inches(0.5), Inches(1.15), "ADMIN", BLUE_MED, admin_steps)

# Petugas flow
petugas_steps = [
    ('oval', 'Login Petugas'),
    ('box', 'Dashboard Petugas'),
    ('diamond', 'Pilih Menu'),
    ('box', 'Kelola Kendaraan'),
    ('box', 'Transaksi Masuk (Input+Barcode)'),
    ('box', 'Transaksi Keluar (Scan+Bayar)'),
    ('box', 'AI Chatbot'),
    ('oval', 'Logout'),
]
draw_role_flow(sl, Inches(4.6), Inches(1.15), "PETUGAS", GREEN, petugas_steps)

# Owner flow
owner_steps = [
    ('oval', 'Login Owner'),
    ('box', 'Dashboard Owner'),
    ('diamond', 'Lihat Rekap?'),
    ('box', 'Input Rentang Tanggal'),
    ('box', 'Grafik + Tabel + Export CSV'),
    ('box', 'AI Chatbot'),
    ('oval', 'Logout'),
]
draw_role_flow(sl, Inches(8.7), Inches(1.15), "OWNER", ORANGE, owner_steps)


# ══════════════════════════════════════════════════
# SLIDE 9: FLOW AI CHATBOT
# ══════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, GRAY_BG)
add_rect(sl, 0, 0, SW, Inches(0.9), fill_color=BLUE_DARK)
add_text_box(sl, Inches(0.8), Inches(0.15), Inches(10), Inches(0.6),
             "07  |  Flow Alur Kerja  -  AI Chatbot", size=24, color=WHITE, bold=True)

fc_bg2 = add_rect(sl, Inches(0.5), Inches(1.1), Inches(12.3), Inches(5.9), fill_color=WHITE, border_color=RGBColor(0xE2,0xE8,0xF0))

# --- Left: Text Chat ---
add_text_box(sl, Inches(0.8), Inches(1.15), Inches(5), Inches(0.35),
             "ALUR CHAT TEKS", size=13, color=BLUE_MED, bold=True)

cx_l = Inches(3.2)
ys = Inches(1.8)
yg = Inches(0.72)

fc_oval(sl, cx_l, ys, Inches(1.3), Inches(0.4), "User Ketik Pesan", fill=BLUE_BG, stroke=BLUE_MED)
fc_arrow_down(sl, cx_l, ys + Inches(0.2), ys + Inches(0.46))

y2c = ys + yg
fc_diamond(sl, cx_l, y2c, Inches(2.0), Inches(0.65), "Deteksi Plat Nomor\natau Kode PKR?", fill=RGBColor(0xFF,0xF7,0xED), stroke=ORANGE)
fc_arrow_down(sl, cx_l, y2c + Inches(0.32), y2c + Inches(0.58))
add_text_box(sl, cx_l + Inches(0.05), y2c + Inches(0.23), Inches(0.4), Inches(0.2), "Ya", size=8, color=GREEN, bold=True)

y3c = y2c + yg + Inches(0.15)
fc_box(sl, cx_l, y3c, Inches(2.0), Inches(0.45), "Query Database\n(Kendaraan + Transaksi Real-time)", fill=GRAY_LIGHT, stroke=GREEN, font_size=8)
fc_arrow_down(sl, cx_l, y3c + Inches(0.22), y3c + Inches(0.48))

y4c = y3c + yg
fc_box(sl, cx_l, y4c, Inches(2.3), Inches(0.45), "Bangun System Prompt\n(Tarif + Area + Statistik + Data Kendaraan)", fill=GRAY_LIGHT, stroke=BLUE_MED, font_size=8)
fc_arrow_down(sl, cx_l, y4c + Inches(0.22), y4c + Inches(0.48))

y5c = y4c + yg
fc_box(sl, cx_l, y5c, Inches(2.0), Inches(0.45), "Panggil NVIDIA API\n(LLaMA 3.1 8B Instruct)", fill=RGBColor(0xEC,0xFD,0xF5), stroke=GREEN, font_size=8)
fc_arrow_down(sl, cx_l, y5c + Inches(0.22), y5c + Inches(0.48))

y6c = y5c + yg
fc_box(sl, cx_l, y6c, Inches(2.2), Inches(0.48), "Tampilkan Respon AI\n+ Parking Card + Barcode", fill=RGBColor(0xDB,0xEA,0xFE), stroke=BLUE_MED, font_size=8)

# --- Right: Image Scan ---
add_text_box(sl, Inches(7.0), Inches(1.15), Inches(5), Inches(0.35),
             "ALUR SCAN GAMBAR", size=13, color=GREEN, bold=True)

cx_r = Inches(9.5)

fc_oval(sl, cx_r, ys, Inches(1.5), Inches(0.4), "Kamera / Upload Foto", fill=BLUE_BG, stroke=GREEN)
fc_arrow_down(sl, cx_r, ys + Inches(0.2), ys + Inches(0.46))

y2r = ys + yg
fc_box(sl, cx_r, y2r, Inches(2.2), Inches(0.45), "Kirim Base64 Image\nke /chatbot/scan-image", fill=GRAY_LIGHT, stroke=GREEN, font_size=8)
fc_arrow_down(sl, cx_r, y2r + Inches(0.22), y2r + Inches(0.48))

y3r = y2r + yg
fc_box(sl, cx_r, y3r, Inches(2.2), Inches(0.45), "NVIDIA Vision API\n(Phi-3.5 Vision) Extract Teks", fill=RGBColor(0xFF,0xF7,0xED), stroke=ORANGE, font_size=8)
fc_arrow_down(sl, cx_r, y3r + Inches(0.22), y3r + Inches(0.48))

y4r = y3r + yg
fc_diamond(sl, cx_r, y4r, Inches(2.0), Inches(0.65), "Terdeteksi PKR\natau Plat?", fill=RGBColor(0xFF,0xF7,0xED), stroke=ORANGE)

# Ya
fc_arrow_down(sl, cx_r, y4r + Inches(0.32), y4r + Inches(0.58))
add_text_box(sl, cx_r + Inches(0.05), y4r + Inches(0.23), Inches(0.4), Inches(0.2), "Ya", size=8, color=GREEN, bold=True)

y5r = y4r + yg + Inches(0.15)
fc_box(sl, cx_r, y5r, Inches(2.0), Inches(0.45), "Query DB + NVIDIA AI\nGenerate Respon Lengkap", fill=RGBColor(0xEC,0xFD,0xF5), stroke=GREEN, font_size=8)
fc_arrow_down(sl, cx_r, y5r + Inches(0.22), y5r + Inches(0.48))

y6r = y5r + yg
fc_box(sl, cx_r, y6r, Inches(2.2), Inches(0.48), "Tampilkan Respon AI\n+ Parking Card + Barcode", fill=RGBColor(0xDB,0xEA,0xFE), stroke=BLUE_MED, font_size=8)

# Divider
div2 = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.3), Inches(1.3), Pt(2), Inches(5.5))
div2.fill.solid()
div2.fill.fore_color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
div2.line.fill.background()


# ══════════════════════════════════════════════════
# SLIDE 10: PENGUJIAN & RENCANA PENGEMBANGAN
# ══════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl, GRAY_BG)
add_rect(sl, 0, 0, SW, Inches(0.9), fill_color=BLUE_DARK)
add_text_box(sl, Inches(0.8), Inches(0.15), Inches(10), Inches(0.6),
             "08  |  Pengujian & Rencana Pengembangan", size=24, color=WHITE, bold=True)

# Left: Pengujian
add_rect(sl, Inches(0.5), Inches(1.15), Inches(6.0), Inches(5.7), fill_color=WHITE, border_color=RGBColor(0xE2,0xE8,0xF0))
bar_test = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.15), Inches(6.0), Inches(0.45))
bar_test.fill.solid()
bar_test.fill.fore_color.rgb = GREEN
bar_test.line.fill.background()
add_text_box(sl, Inches(0.7), Inches(1.17), Inches(5.5), Inches(0.4),
             "HASIL PENGUJIAN (Black-Box Testing)", size=14, color=WHITE, bold=True)

test_items = [
    ("Autentikasi & Otorisasi", True, GREEN),
    "  Login berhasil membedakan sesi per role (admin/petugas/owner)",
    "  Akun nonaktif ditolak saat login",
    "",
    ("Transaksi Parkir (Masuk & Keluar)", True, GREEN),
    "  Kode PKR auto-generate + barcode, perhitungan biaya akurat",
    "  Area penuh ditolak, duplikasi kendaraan masuk dicegah",
    "",
    ("AI Chatbot", True, GREEN),
    "  Deteksi plat nomor & kode PKR dari teks dan gambar",
    "  Parking Card + barcode ditampilkan dengan benar",
    "",
    ("Validasi Input", True, GREEN),
    "  Username/plat duplikat ditolak, kolom wajib tervalidasi",
    "",
    ("Kesimpulan: LULUS uji fungsional standar", True, BLUE_MED),
]
add_multiline(sl, Inches(0.7), Inches(1.75), Inches(5.5), Inches(5.0), test_items, size=10, color=GRAY_DARK)

# Right: Rencana
add_rect(sl, Inches(6.8), Inches(1.15), Inches(6.0), Inches(5.7), fill_color=WHITE, border_color=RGBColor(0xE2,0xE8,0xF0))
bar_dev = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.15), Inches(6.0), Inches(0.45))
bar_dev.fill.solid()
bar_dev.fill.fore_color.rgb = ORANGE
bar_dev.line.fill.background()
add_text_box(sl, Inches(7.0), Inches(1.17), Inches(5.5), Inches(0.4),
             "RENCANA PENGEMBANGAN", size=14, color=WHITE, bold=True)

dev_items = [
    ("1. Real-time Notification", True, ORANGE),
    "   Notifikasi otomatis via Email/WhatsApp/Push",
    "   saat ada event penting (area penuh, anomali, dll)",
    "",
    ("2. Export Laporan PDF/Excel", True, ORANGE),
    "   Ekspor rekapitulasi ke PDF dan Excel untuk evaluasi",
    "",
    ("3. Multi-turn AI Chatbot", True, ORANGE),
    "   Riwayat percakapan per sesi + model AI lebih advanced",
    "",
    ("4. Pembayaran Digital (QRIS)", True, ORANGE),
    "   Integrasi gateway pembayaran cashless",
    "",
    ("5. Deteksi Plat Otomatis (ANPR)", True, ORANGE),
    "   Kamera CCTV + AI computer vision",
    "",
    ("6. PWA / Mobile App", True, ORANGE),
    "   Progressive Web App atau Android/iOS native",
    "",
    ("7. Membership & Langganan", True, ORANGE),
    "   Tarif khusus harian/mingguan/bulanan",
]
add_multiline(sl, Inches(7.0), Inches(1.75), Inches(5.5), Inches(5.0), dev_items, size=10, color=GRAY_DARK)


# ══════════════════════════════════════════════════
# SLIDE 11: PENUTUP / TERIMA KASIH
# ══════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_rect(sl, 0, 0, SW, SH, BLUE_DARK, RGBColor(0x0F, 0x25, 0x40))

deco3 = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.5), Inches(3.5), Inches(6), Inches(6))
deco3.fill.solid()
deco3.fill.fore_color.rgb = RGBColor(0x2B, 0x4F, 0x7F)
deco3.line.fill.background()

deco4 = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(-2), Inches(5), Inches(5))
deco4.fill.solid()
deco4.fill.fore_color.rgb = RGBColor(0x1A, 0x30, 0x50)
deco4.line.fill.background()

add_text_box(sl, Inches(1.5), Inches(1.8), Inches(10), Inches(1.0),
             "Terima Kasih", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

line2 = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.0), Inches(2.3), Pt(3))
line2.fill.solid()
line2.fill.fore_color.rgb = BLUE_LIGHT
line2.line.fill.background()

add_text_box(sl, Inches(1.5), Inches(3.4), Inches(10), Inches(0.8),
             "SmartPark  -  Sistem Manajemen Parkir Terintegrasi AI",
             size=18, color=RGBColor(0xA0, 0xC4, 0xF0), align=PP_ALIGN.CENTER)

add_text_box(sl, Inches(1.5), Inches(4.3), Inches(10), Inches(1.0),
             "Laravel 12  |  React 19  |  Vite  |  Tailwind CSS  |  NVIDIA AI API\nMySQL  |  Sanctum Authentication  |  Barcode System",
             size=14, color=RGBColor(0x80, 0xA0, 0xC0), align=PP_ALIGN.CENTER)

add_text_box(sl, Inches(1.5), Inches(5.5), Inches(10), Inches(0.6),
             "github.com/HumamasyariDev/humanityalok",
             size=13, color=RGBColor(0x70, 0x90, 0xB0), align=PP_ALIGN.CENTER)

add_text_box(sl, Inches(1.5), Inches(6.2), Inches(10), Inches(0.5),
             "HumamasyariDev  -  2026",
             size=13, color=RGBColor(0x60, 0x80, 0xA0), align=PP_ALIGN.CENTER)


# ── Save ──
output_path = r"C:\laragon\www\humanityalok\SmartPark_Presentasi.pptx"
prs.save(output_path)
print(f"PPT berhasil dibuat: {output_path}")
print(f"Total slides: {len(prs.slides)}")
